# import pptx
# from data_extractor.file_loaders.file_loader import FileLoader

# class PPTLoader(FileLoader):
#     def validate_file(self, file_path: str) -> bool:
#         return file_path.lower().endswith('.pptx') or file_path.lower().endswith('.ppt')

#     def load_file(self, file_path: str) -> pptx.Presentation:
#         if not self.validate_file(file_path):
#             raise ValueError("Invalid PPT file.")
        
#         try:
#             # Attempt to load the PPTX file
#             return pptx.Presentation(file_path)
#         except Exception:
#             # Catch any exception related to loading the file and raise the expected error
#             raise ValueError("Invalid PPT file.")



import pptx
import msoffcrypto
from data_extractor.file_loaders.file_loader import FileLoader
from io import BytesIO

class PPTLoader(FileLoader):
    def validate_file(self, file_path: str) -> bool:
        return file_path.lower().endswith('.pptx') or file_path.lower().endswith('.ppt')

    def load_file(self, file_path: str) -> pptx.Presentation:
        if not self.validate_file(file_path):
            raise ValueError("Invalid PPT file.")
        
        try:
            # Attempt to load the PPTX file
            presentation = None

            with open(file_path, 'rb') as f:
                # Create a msoffcrypto OfficeFile object
                office_file = msoffcrypto.OfficeFile(f)

                # Check if the file is encrypted
                if office_file.is_encrypted:
                    password = str(input("Enter your PPT password: "))  # Prompt user for password
                    office_file.load_key(password=password)  # Load the provided password
                    decrypted_file = BytesIO()
                    office_file.decrypt(decrypted_file)  # Decrypt the file
                    decrypted_file.seek(0)  # Reset the BytesIO cursor to the beginning
                    presentation = pptx.Presentation(decrypted_file)  # Load the decrypted presentation
                else:
                    # If not encrypted, load it directly
                    presentation = pptx.Presentation(f)
            return presentation

        except msoffcrypto.exceptions.InvalidKeyError:
            raise ValueError("Failed to decrypt the PPT file: Incorrect password.")
        except ValueError as ve:
            raise ve
        except Exception as e:
            raise ValueError(f"Invalid PPT file: {str(e)}")
